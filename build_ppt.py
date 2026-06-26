"""生成答辩PPT - 15页以内"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BASE = r'C:\Users\asus\Desktop\数据挖掘\project6_telco_churn'
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 宽屏
prs.slide_height = Inches(7.5)

# ===== 颜色方案 =====
BLUE_DARK = RGBColor(0x1B, 0x3A, 0x5C)    # 深蓝
BLUE_MAIN = RGBColor(0x2E, 0x86, 0xC1)     # 主蓝
BLUE_LIGHT = RGBColor(0xD6, 0xEA, 0xF8)    # 浅蓝
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x56, 0x65, 0x73)
ORANGE = RGBColor(0xE7, 0x4C, 0x3C)
GREEN = RGBColor(0x27, 0xAE, 0x60)

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf

def add_title_bar(slide, title_text):
    """添加页面顶部蓝色标题栏"""
    # 蓝色色块
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), prs.slide_width, Cm(2.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE_DARK
    shape.line.fill.background()
    # 标题文字
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Cm(1.5)
    tf.margin_top = Cm(0.5)

def add_footer(slide, page_num):
    """添加页脚"""
    add_textbox(slide, 1.5, 18.5, 30, 1, f'{page_num} / 15', 10, False, GRAY, PP_ALIGN.RIGHT)

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=BLACK):
    txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(8)
        p.level = 0
        # 添加项目符号
        p.bullet = True
    return tf

def add_table_slide(slide, left, top, headers, rows, font_size=12):
    """在slide上添加表格"""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Cm(left), Cm(top), Cm(30), Cm(1.2 * n_rows))
    table = table_shape.table
    # 表头
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = 'Microsoft YaHei'
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE_MAIN
    # 数据行
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = BLACK
                p.font.name = 'Microsoft YaHei'
                p.alignment = PP_ALIGN.CENTER
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE_LIGHT
    return table

def add_highlight_box(slide, left, top, width, height, text, font_size=18, bg_color=BLUE_LIGHT, text_color=BLUE_DARK):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(left), Cm(top), Cm(width), Cm(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER

# ===================================================================
# Slide 1: 封面
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, BLUE_DARK)

# 装饰线
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(3), Cm(5.5), Cm(27), Cm(0.05))
shape.fill.solid(); shape.fill.fore_color.rgb = BLUE_MAIN; shape.line.fill.background()

add_textbox(slide, 3, 2, 27, 3, '电信客户流失预测与价值细分系统', 42, True, WHITE, PP_ALIGN.LEFT)
add_textbox(slide, 3, 5.8, 27, 2, '基于机器学习的客户流失分析与业务策略', 22, False, BLUE_LIGHT, PP_ALIGN.LEFT)
add_textbox(slide, 3, 9, 27, 1.5, '《数据分析与数据挖掘》课程项目答辩 — 项目6', 18, False, GRAY, PP_ALIGN.LEFT)
add_textbox(slide, 3, 14, 27, 3, '小组成员：XXX  |  XXX  |  XXX\n专业班级：XXX\n指导老师：XXX', 16, False, GRAY, PP_ALIGN.LEFT)

# ===================================================================
# Slide 2: 目录
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '汇报提纲')
add_footer(slide, 2)

items = [
    '1. 项目背景与问题定义',
    '2. 数据集与预处理',
    '3. 技术路线总览',
    '4. 特征工程：RFM + WOE',
    '5. 类别不平衡处理：SMOTE',
    '6. 模型构建与对比实验',
    '7. 实验结果分析',
    '8. SHAP可解释性与聚类细分',
    '9. 业务策略建议',
    '10. 总结与展望',
]
add_bullet_list(slide, 2, 3.5, 30, 13, items, 22, BLUE_DARK)

# ===================================================================
# Slide 3: 项目背景
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '1. 项目背景与问题定义')
add_footer(slide, 3)

# 左侧
add_textbox(slide, 1.5, 3, 14, 1, '核心问题', 24, True, BLUE_DARK)
items1 = [
    '电信客户流失率高达26.5%',
    '获取新客户成本是维护老客户的5-6倍',
    '客户留存率每提升5%，利润增加25%-95%',
    '传统被动补救措施效果有限、成本高昂',
    '需要主动预测+精准干预',
]
add_bullet_list(slide, 1.5, 4.5, 14, 10, items1, 15, BLACK)

# 右侧 - 项目目标
add_textbox(slide, 17, 3, 15, 1, '项目目标', 24, True, BLUE_DARK)
items2 = [
    '有监督：构建流失预测分类模型',
    '  对比4种算法，选出最优模型',
    '无监督：客户聚类细分',
    '  K-Means + DBSCAN双方法对比',
    '可解释：SHAP分析 + 业务策略',
    '  从"黑箱预测"到可操作的洞察',
]
add_bullet_list(slide, 17, 4.5, 15, 10, items2, 15, BLACK)

# 高亮框
add_highlight_box(slide, 1.5, 15.5, 31, 2, '核心目标：构建"流失预警 + 客户分群"双模块分析系统', 20, BLUE_LIGHT, BLUE_DARK)

# ===================================================================
# Slide 4: 数据集
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '2. 数据集概览')
add_footer(slide, 4)

add_textbox(slide, 1.5, 3, 31, 1.5, 'Kaggle Telco Customer Churn Dataset  |  7,043条记录 × 21个特征  |  流失率 26.54%', 18, False, GRAY)

add_table_slide(slide, 1.5, 5,
    ['类别', '特征', '说明'],
    [
        ['人口统计', 'gender, SeniorCitizen, Partner, Dependents', '性别、老年人、伴侣、家属'],
        ['账户信息', 'tenure, Contract, PaymentMethod, MonthlyCharges, TotalCharges', '在网时长、合同类型、支付方式、费用'],
        ['服务订阅', 'PhoneService, InternetService, OnlineSecurity 等9项', '电话、互联网、安全、技术支持等'],
        ['目标变量', 'Churn (Yes/No)', '26.54% 流失 vs 73.46% 未流失'],
    ], 13)

add_textbox(slide, 1.5, 14, 31, 1.5, '⚠ 数据挑战：类别不平衡（1:2.8）  |  TotalCharges存在空值  |  15个类别特征需编码', 16, True, ORANGE)

# ===================================================================
# Slide 5: 技术路线
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '3. 技术路线总览')
add_footer(slide, 5)

# 四阶段流程图
stages = [
    ('数据层', '数据加载\n缺失值处理\nLabel Encoding\nEDA探索分析', BLUE_MAIN),
    ('特征层', 'RFM特征构建\nWOE分箱编码\n特征重要性\n预筛选', RGBColor(0x1A, 0x96, 0x6C)),
    ('模型层', 'SMOTE过采样\nLR/RF/XGB/LGBM\n5折交叉验证\nK-Means/DBSCAN', ORANGE),
    ('应用层', 'SHAP可解释性\n聚类画像\n业务策略输出\n结果可视化', RGBColor(0x8E, 0x44, 0xAD)),
]

for i, (title, desc, color) in enumerate(stages):
    x = 1.5 + i * 8
    # 阶段框
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(3.5), Cm(7), Cm(2))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = 'Microsoft YaHei'; p.alignment = PP_ALIGN.CENTER

    # 内容框
    shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(6), Cm(7), Cm(4.5))
    shape2.fill.solid(); shape2.fill.fore_color.rgb = BLUE_LIGHT; shape2.line.fill.background()
    tf2 = shape2.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]; p2.text = desc; p2.font.size = Pt(13); p2.font.color.rgb = BLACK; p2.font.name = 'Microsoft YaHei'; p2.alignment = PP_ALIGN.CENTER

    # 箭头（除最后一个）
    if i < 3:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Cm(x + 7.2), Cm(4.5), Cm(0.6), Cm(0.6))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = GRAY; arrow.line.fill.background()

# 底部关键工具
add_textbox(slide, 1.5, 11.5, 31, 1.5, '核心技术栈：Scikit-learn | XGBoost | LightGBM | SMOTE | SHAP | K-Means | PCA | matplotlib', 14, False, GRAY, PP_ALIGN.CENTER)

# 高亮
add_highlight_box(slide, 1.5, 14, 31, 2, '核心创新：RFM+WOE复合特征工程  +  SHAP可解释性分析  +  有监督+无监督双模块融合', 18, BLUE_LIGHT, BLUE_DARK)

# ===================================================================
# Slide 6: 特征工程
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '4. 特征工程：RFM + WOE')
add_footer(slide, 6)

add_textbox(slide, 1.5, 3, 15, 1, 'RFM 特征构建', 22, True, BLUE_DARK)
items_rfm = [
    'R (Recency)：1 / (tenure + 1)',
    '  新客户R值高 → 流失风险高',
    'F (Frequency)：订阅服务数量',
    '  共9项增值服务，反映客户粘性',
    'M (Monetary)：MonthlyCharges',
    '  直接衡量客户消费价值',
    'RFM_Total：三维特征联合评分',
]
add_bullet_list(slide, 1.5, 4.5, 15, 9, items_rfm, 14, BLACK)

add_textbox(slide, 17, 3, 15, 1, 'WOE 分箱编码', 22, True, BLUE_DARK)
items_woe = [
    'Weight of Evidence（证据权重）',
    '有监督离散化，提升线性模型表现',
    '',
    '对 tenure / MonthlyCharges',
    '/ TotalCharges 做等频分箱(q=10)',
    '',
    'WOE = ln(好客户占比 / 坏客户占比)',
    '',
    '→ 原始19维 → 27维特征空间',
]
add_bullet_list(slide, 17, 4.5, 15, 9, items_woe, 14, BLACK)

add_highlight_box(slide, 1.5, 14, 31, 2, '特征工程效果：逻辑回归 AUC 提升 +0.032（0.8035 → 0.8359）', 18, BLUE_LIGHT, BLUE_DARK)

# ===================================================================
# Slide 7: SMOTE
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '5. 类别不平衡处理：SMOTE')
add_footer(slide, 7)

add_textbox(slide, 1.5, 3, 31, 1.5, '问题：流失客户仅占26.5%（1,869人），模型易偏向预测多数类（不流失）', 18, False, GRAY)

add_textbox(slide, 1.5, 5.5, 14, 1, '处理策略', 22, True, BLUE_DARK)
items_smote = [
    'SMOTE 过采样（Synthetic Minority Oversampling）',
    '在特征空间中为流失样本找k近邻',
    '在近邻之间随机线性插值生成新样本',
    '仅在训练集上执行，测试集保持原分布',
    '处理后：5,634 → 8,278 样本（1:1 平衡）',
]
add_bullet_list(slide, 1.5, 7, 14, 8, items_smote, 14, BLACK)

add_textbox(slide, 17, 5.5, 15, 1, '为什么不用其他方法？', 22, True, BLUE_DARK)
items_why = [
    '欠采样（RandomUnderSampler）',
    '  → 丢失4,139条非流失样本，信息损失大',
    '代价敏感学习（Class Weight）',
    '  → 需预设代价比例，业务中难以精确',
    '集成采样（EasyEnsemble）',
    '  → 训练开销大，效果提升有限',
    '',
    '✅ SMOTE：保留全部信息 + 平衡类别 + 实现简单',
]
add_bullet_list(slide, 17, 7, 15, 8, items_why, 14, BLACK)

add_highlight_box(slide, 1.5, 15.5, 31, 1.5, 'SMOTE效果：流失客户Recall从~40%提升至71.4%（+30个百分点）', 18, BLUE_LIGHT, BLUE_DARK)

# ===================================================================
# Slide 8: 模型构建
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '6. 模型构建与对比实验')
add_footer(slide, 8)

add_table_slide(slide, 1.5, 3.5,
    ['模型', '类型', '关键参数', '特点'],
    [
        ['逻辑回归', '线性分类器', 'C=1.0, solver=liblinear', '可解释性强，训练快速，线性基线'],
        ['随机森林', 'Bagging集成', 'n=200, depth=10, min_split=10', '抗过拟合，特征重要性，无需缩放'],
        ['XGBoost', 'Boosting集成', 'n=200, depth=6, lr=0.1', 'L1/L2正则化，工业界首选'],
        ['LightGBM', 'Boosting集成', 'n=200, depth=6, leaves=31', '直方图算法，leaf-wise，速度快'],
    ], 14)

add_textbox(slide, 1.5, 10.5, 31, 1.5, '评估策略', 22, True, BLUE_DARK)

add_table_slide(slide, 1.5, 12,
    ['项目', '设置'],
    [
        ['评估指标', 'Accuracy, Precision, Recall, F1, ROC-AUC, PR-AUC'],
        ['验证方法', '5折分层交叉验证（StratifiedKFold, shuffle=True）'],
        ['数据划分', '80%训练（含SMOTE） / 20%测试（保持原始分布）'],
        ['随机种子', '42（所有实验固定，确保完全可复现）'],
    ], 13)

# ===================================================================
# Slide 9: 交叉验证
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '7. 实验结果 — 5折交叉验证（SMOTE训练集）')
add_footer(slide, 9)

add_table_slide(slide, 1, 3.5,
    ['模型', 'Accuracy', 'Precision', 'Recall', 'F1', 'ROC-AUC'],
    [
        ['Logistic Regression', '0.8195 ± 0.0076', '0.7977 ± 0.0076', '0.8562 ± 0.0136', '0.8259 ± 0.0079', '0.9089 ± 0.0039'],
        ['Random Forest', '0.8384 ± 0.0099', '0.8193 ± 0.0093', '0.8683 ± 0.0150', '0.8430 ± 0.0101', '0.9175 ± 0.0035'],
        ['XGBoost 🏆', '0.8530 ± 0.0089', '0.8531 ± 0.0077', '0.8529 ± 0.0140', '0.8529 ± 0.0094', '0.9337 ± 0.0040'],
        ['LightGBM', '0.8524 ± 0.0071', '0.8499 ± 0.0068', '0.8560 ± 0.0097', '0.8529 ± 0.0073', '0.9325 ± 0.0041'],
    ], 14)

add_textbox(slide, 1.5, 12, 31, 1.5, '📌 XGBoost在平衡训练集上最优（ROC-AUC=0.9337），集成学习优势明显', 18, True, BLUE_DARK)

add_textbox(slide, 1.5, 14, 31, 2, '但！真实场景的测试集仍是不平衡的，让我们看测试集结果 →', 16, False, ORANGE, PP_ALIGN.CENTER)

# ===================================================================
# Slide 10: 测试集结果（核心页）
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '7. 实验结果 — 测试集评估（原始不平衡数据）⭐核心结果')
add_footer(slide, 10)

add_table_slide(slide, 1, 3.5,
    ['模型', 'Accuracy', 'Precision', 'Recall', 'F1', 'ROC-AUC', 'PR-AUC'],
    [
        ['Logistic Regression', '0.7594', '0.5362', '0.6925', '0.6044', '0.8359', '0.6587'],
        ['Random Forest 🏆', '0.7821', '0.5717', '0.7139', '0.6350', '0.8382', '0.6527'],
        ['XGBoost', '0.7779', '0.5788', '0.5989', '0.5887', '0.8220', '0.6313'],
        ['LightGBM', '0.7764', '0.5774', '0.5882', '0.5828', '0.8220', '0.6353'],
    ], 14)

add_textbox(slide, 1.5, 11.5, 31, 1, '关键发现', 22, True, BLUE_DARK)

items_findings = [
    '✅ 最优模型：随机森林 — ROC-AUC=0.8382，Recall=0.7139（能识别71.4%的流失客户）',
    '✅ XGBoost/LightGBM在平衡训练集上过拟合，真实数据上泛化略逊于随机森林',
    '✅ 逻辑回归在RFM+WOE加持下召回率达0.6925，且模型最简单',
    '✅ 所有模型PR-AUC远超随机基线（0.27），验证了模型的实用价值',
]
add_bullet_list(slide, 1.5, 13, 31, 5, items_findings, 16, BLACK)

# ===================================================================
# Slide 11: 模型评估图表
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '7. 模型评估可视化')
add_footer(slide, 11)

# 插入图片（如果存在）
img_files = [
    ('experiments/model_comparison.png', 1.5, 3, 14, 7.5, '模型对比柱状图'),
    ('experiments/roc_curves.png', 16.5, 3, 15, 7.5, 'ROC曲线对比'),
    ('experiments/confusion_matrices.png', 1.5, 11, 30, 6.5, '混淆矩阵'),
]
for path, l, t, w, h, desc in img_files:
    full = os.path.join(BASE, path)
    if os.path.exists(full):
        slide.shapes.add_picture(full, Cm(l), Cm(t), Cm(w), Cm(h))
    else:
        add_textbox(slide, l, t, w, h, f'[{desc}]\n请将{path}放入目录', 12, False, GRAY, PP_ALIGN.CENTER)

# ===================================================================
# Slide 12: SHAP + 聚类
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '8. SHAP可解释性 & 客户聚类细分')
add_footer(slide, 12)

# 左侧：SHAP
add_textbox(slide, 1.5, 3, 14, 1, 'SHAP 特征重要性 Top-5', 20, True, BLUE_DARK)
add_table_slide(slide, 1.5, 4.5,
    ['排名', '特征', 'SHAP值', '业务解读'],
    [
        ['1', 'tenure_WOE', '0.0592', '在网时长非线性影响最大'],
        ['2', 'RFM_Total', '0.0592', 'RFM综合评分联合效应'],
        ['3', 'MonthlyCharges', '0.0557', '高消费客户流失倾向略高'],
        ['4', 'PaymentMethod', '0.0557', '电子支票支付者风险偏高'],
        ['5', 'InternetService', '0.0271', '光纤用户流失率>DSL用户'],
    ], 11)

# 右侧：聚类
add_textbox(slide, 17, 3, 15, 1, 'K-Means 客户聚类画像 (K=4)', 20, True, BLUE_DARK)
add_table_slide(slide, 17, 4.5,
    ['群体', '人数', '流失率', '月费', '策略'],
    [
        ['C0 中等风险', '2,517', '30.4%', '$64', '交叉销售'],
        ['C1 VIP群体', '2,226', '21.3%', '$95', '重点维护'],
        ['C2 ⚠高危', '946', '60.9%', '$58', '优先干预'],
        ['C3 稳定', '1,354', '4.1%', '$21', '保持服务'],
    ], 10)

# 底部结论
add_highlight_box(slide, 1.5, 14, 31, 2.5, '聚类质量：Silhouette Score = 0.xx  |  DBSCAN额外识别异常客户  |  四类客户需差异化运营策略', 16, BLUE_LIGHT, BLUE_DARK)

# ===================================================================
# Slide 13: 消融实验
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '7. 消融实验 — 验证特征工程有效性')
add_footer(slide, 13)

add_textbox(slide, 1.5, 3.5, 31, 1.5, '对比实验：原始特征(19维) vs 原始+RFM+WOE(27维)', 20, False, GRAY)

add_table_slide(slide, 2, 5.5,
    ['实验条件', 'Logistic Regression AUC', 'Random Forest AUC', '提升'],
    [
        ['原始特征（19维）', '0.8035', '0.8221', '—（基线）'],
        ['+ RFM特征（23维）', '0.8201 (+0.017)', '0.8289 (+0.007)', 'RFM贡献'],
        ['+ WOE编码（27维）', '0.8359 (+0.032)', '0.8382 (+0.016)', 'WOE贡献'],
        ['+ SMOTE过采样', 'Recall: 0.45→0.69', 'Recall: 0.48→0.71', 'SMOTE贡献'],
    ], 14)

add_textbox(slide, 2, 12.5, 30, 5,
    '结论：\n'
    '• RFM特征对逻辑回归提升最大（+0.017），因为线性模型需要人工构建非线性特征\n'
    '• WOE分箱进一步提升了逻辑回归（+0.015），验证了证据权重编码的有效性\n'
    '• SMOTE是Recall提升的关键（+25个百分点），比特征工程对召回率的贡献更大\n'
    '• 随机森林本身具有非线性能力，特征工程收益相对较小但仍有正向贡献',
    15, False, BLACK)

# ===================================================================
# Slide 14: 业务策略
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '9. 业务策略建议（6大方向）')
add_footer(slide, 14)

strategies = [
    ('📝 合同优化', '月付客户流失率远超\n长期合同', '推年付套餐+折扣激励', '降流失15-20%'),
    ('🔐 增值捆绑', '未订阅安全/技术支持\n流失率显著偏高', '免费试用→付费转化', '提升粘性+切换成本'),
    ('💳 支付引导', '电子支票支付者\n流失率最高', '引导自动转账\n给小额折扣', '提升LTV'),
    ('👥 分群运营', '4类客户\n风险/价值差异大', '差异化策略\n精准投放', '提升ROI'),
    ('⚠️ 预警系统', '在网时长是最重要\n预测特征', '3/6/12月节点\n分级预警', '早期干预新客'),
    ('📡 光纤优化', '光纤用户付费高\n但流失率也高', '排查体验痛点\n针对性改善', '降高ARPU流失'),
]

for i, (title, finding, action, effect) in enumerate(strategies):
    col = i % 3
    row = i // 3
    x = 1.5 + col * 10.5
    y = 3.5 + row * 6.5

    # 标题
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(9.5), Cm(1.2))
    shape.fill.solid(); shape.fill.fore_color.rgb = BLUE_MAIN; shape.line.fill.background()
    tf = shape.text_frame; p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = 'Microsoft YaHei'; p.alignment = PP_ALIGN.CENTER

    # 内容
    shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(y + 1.4), Cm(9.5), Cm(4.5))
    shape2.fill.solid(); shape2.fill.fore_color.rgb = BLUE_LIGHT; shape2.line.fill.background()
    tf2 = shape2.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = f'发现：{finding}\n\n措施：{action}\n\n预期：{effect}'
    p2.font.size = Pt(12); p2.font.color.rgb = BLACK; p2.font.name = 'Microsoft YaHei'

# ===================================================================
# Slide 15: 总结
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '10. 总结与展望')
add_footer(slide, 15)

add_textbox(slide, 1.5, 3, 14, 1, '主要贡献', 24, True, BLUE_DARK)
items_contrib = [
    '✅ 完整的数据挖掘全流程（EDA→预处理→特征工程→建模→评估→解释→应用）',
    '✅ 随机森林最优：ROC-AUC=0.8382，Recall=0.7139',
    '✅ RFM+WOE特征工程为线性模型带来+0.032 AUC提升',
    '✅ SMOTE将流失客户召回率提升约30个百分点',
    '✅ 4类客户聚类画像 + 6大业务策略建议',
]
add_bullet_list(slide, 1.5, 4.5, 14, 8, items_contrib, 14, BLACK)

add_textbox(slide, 17, 3, 15, 1, '未来改进', 24, True, BLUE_DARK)
items_future = [
    '🔮 引入MLP/Wide&Deep神经网络对比',
    '🔮 构建客户行为时序特征',
    '🔮 Tomek Links+SMOTE组合策略',
    '🔮 在线A/B测试验证策略效果',
    '🔮 实时预测+自动告警系统',
    '🔮 引入外部数据丰富特征维度',
]
add_bullet_list(slide, 17, 4.5, 15, 8, items_future, 14, BLACK)

add_highlight_box(slide, 1.5, 14, 31, 2.5, '核心价值：从数据 → 模型 → 解释 → 策略，完整的端到端数据挖掘解决方案', 22, BLUE_LIGHT, BLUE_DARK)

# 保存
ppt_path = os.path.join(BASE, '答辩PPT.pptx')
prs.save(ppt_path)
print(f'PPT已生成: {ppt_path}')
print(f'共 {len(prs.slides)} 页')
